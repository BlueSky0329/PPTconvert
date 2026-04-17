import os
import tempfile
import unittest

from pptx import Presentation

from core.template_manager import TemplateManager
from core.template_style import delete_all_slides


class TemplateManagerTest(unittest.TestCase):
    def test_load_template_recreates_presentation_from_cached_bytes(self):
        tmp_path = None
        try:
            presentation = Presentation()
            presentation.slides.add_slide(presentation.slide_layouts[6])
            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
                tmp_path = tmp.name
            presentation.save(tmp_path)

            manager = TemplateManager()
            first = manager.load_template(tmp_path)
            delete_all_slides(first)
            self.assertEqual(len(first.slides), 0)

            second = manager.load_template(tmp_path)
            self.assertEqual(len(second.slides), 1)
        finally:
            if tmp_path and os.path.exists(tmp_path):
                os.unlink(tmp_path)


if __name__ == "__main__":
    unittest.main()
