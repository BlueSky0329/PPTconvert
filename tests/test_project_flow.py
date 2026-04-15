import os
import tempfile
import unittest
from unittest.mock import patch

from docx import Document
from pptx import Presentation

from core import pdf_exam_extract
from domain.models import ExamProject
from exporters.manifest_json import load_project_manifest_project
from workflows.project_flow import build_pdf_project, build_word_project, export_project_outputs


class ProjectFlowTest(unittest.TestCase):
    def _write_minimal_word_exam(self, path: str) -> None:
        document = Document()
        document.add_paragraph("四. 数量关系")
        document.add_paragraph("66. 甲、乙两队合修一段公路，若甲单独修需要12天，乙单独修需要18天，两队合修几天完成？")
        document.add_paragraph("A. 6")
        document.add_paragraph("B. 7")
        document.add_paragraph("C. 8")
        document.add_paragraph("D. 9")
        document.save(path)

    def _write_minimal_pdf_exam(self, path: str) -> None:
        document = pdf_exam_extract.fitz.open()
        try:
            page = document.new_page(width=595, height=842)
            lines = [
                "66. How many hours will the two teams need together?",
                "A. 6",
                "B. 7",
                "C. 8",
                "D. 9",
            ]
            y = 72
            for line in lines:
                page.insert_text((72, y), line, fontsize=14)
                y += 24
            document.save(path)
        finally:
            document.close()

    def test_build_pdf_project_passes_mode_through_to_pdf_builder(self):
        project = ExamProject(title="示例工程")

        with patch(
            "workflows.project_flow.build_exam_project_from_pdf",
            return_value=project,
        ) as build_mock, patch(
            "workflows.project_flow.select_project",
            side_effect=lambda current_project, **_kwargs: current_project,
        ) as select_mock, patch(
            "workflows.project_flow.annotate_project_quality",
        ):
            result_project, asset_dir = build_pdf_project(
                "sample.pdf",
                mode="data",
                question_range_spec="111-115",
                asset_dir="assets",
            )

        self.assertIs(result_project, project)
        self.assertEqual(asset_dir, "assets")
        build_mock.assert_called_once_with(
            "sample.pdf",
            mode="data",
            asset_dir="assets",
            document_subject_hint=None,
        )
        self.assertEqual(select_mock.call_args.kwargs["subjects"], ["data"])
        self.assertEqual(
            [(item.start, item.end) for item in select_mock.call_args.kwargs["question_ranges"]],
            [(111, 115)],
        )

    def test_word_project_flow_exports_real_files(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            docx_input = os.path.join(temp_dir, "input.docx")
            asset_dir = os.path.join(temp_dir, "assets")
            docx_output = os.path.join(temp_dir, "booklet.docx")
            pptx_output = os.path.join(temp_dir, "slides.pptx")
            manifest_output = os.path.join(temp_dir, "project.json")
            self._write_minimal_word_exam(docx_input)

            project, questions, chosen_asset_dir = build_word_project(
                docx_input,
                asset_dir=asset_dir,
            )
            outputs = export_project_outputs(
                project,
                asset_dir=chosen_asset_dir,
                docx_output=docx_output,
                ppt_output=pptx_output,
                manifest_output=manifest_output,
            )

            self.assertEqual(len(questions), 1)
            self.assertEqual(project.question_count, 1)
            self.assertEqual(project.sections[0].kind, "quant")
            self.assertTrue(os.path.isdir(chosen_asset_dir))
            self.assertEqual(outputs.docx_path, docx_output)
            self.assertEqual(outputs.pptx_path, pptx_output)
            self.assertEqual(outputs.manifest_path, manifest_output)
            self.assertTrue(os.path.isfile(docx_output))
            self.assertTrue(os.path.isfile(pptx_output))
            self.assertTrue(os.path.isfile(manifest_output))

            exported_docx = Document(docx_output)
            paragraph_texts = [paragraph.text for paragraph in exported_docx.paragraphs]
            self.assertTrue(any("数量关系" in text for text in paragraph_texts))
            self.assertIn(
                "66. 甲、乙两队合修一段公路，若甲单独修需要12天，乙单独修需要18天，两队合修几天完成？",
                paragraph_texts,
            )

            exported_pptx = Presentation(pptx_output)
            self.assertEqual(len(exported_pptx.slides), 1)

            manifest_project = load_project_manifest_project(manifest_output)
            self.assertEqual(manifest_project.question_count, 1)
            self.assertEqual(manifest_project.sections[0].kind, "quant")
            self.assertEqual(manifest_project.sections[0].questions[0].source_number, "66")

    @unittest.skipIf(pdf_exam_extract.fitz is None, "PyMuPDF not installed")
    def test_pdf_project_flow_exports_real_files(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            pdf_input = os.path.join(temp_dir, "input.pdf")
            asset_dir = os.path.join(temp_dir, "assets")
            docx_output = os.path.join(temp_dir, "booklet.docx")
            pptx_output = os.path.join(temp_dir, "slides.pptx")
            manifest_output = os.path.join(temp_dir, "project.json")
            self._write_minimal_pdf_exam(pdf_input)

            project, chosen_asset_dir = build_pdf_project(
                pdf_input,
                mode="quant",
                asset_dir=asset_dir,
                document_subject_hint="quant",
            )
            outputs = export_project_outputs(
                project,
                asset_dir=chosen_asset_dir,
                docx_output=docx_output,
                ppt_output=pptx_output,
                manifest_output=manifest_output,
            )

            self.assertEqual(project.question_count, 1)
            self.assertEqual(project.sections[0].kind, "quant")
            self.assertEqual(project.sections[0].questions[0].source_number, "66")
            self.assertIn("How many hours", project.sections[0].questions[0].stem)
            self.assertTrue(os.path.isdir(chosen_asset_dir))
            self.assertEqual(outputs.docx_path, docx_output)
            self.assertEqual(outputs.pptx_path, pptx_output)
            self.assertEqual(outputs.manifest_path, manifest_output)
            self.assertTrue(os.path.isfile(docx_output))
            self.assertTrue(os.path.isfile(pptx_output))
            self.assertTrue(os.path.isfile(manifest_output))

            exported_docx = Document(docx_output)
            paragraph_texts = [paragraph.text for paragraph in exported_docx.paragraphs]
            self.assertTrue(any("数量关系" in text for text in paragraph_texts))
            self.assertTrue(any("How many hours" in text for text in paragraph_texts))

            exported_pptx = Presentation(pptx_output)
            self.assertEqual(len(exported_pptx.slides), 1)

            manifest_project = load_project_manifest_project(manifest_output)
            self.assertEqual(manifest_project.question_count, 1)
            self.assertEqual(manifest_project.sections[0].kind, "quant")
            self.assertEqual(manifest_project.sections[0].questions[0].source_number, "66")

    @unittest.skipIf(pdf_exam_extract.fitz is not None, "PyMuPDF installed")
    def test_build_pdf_project_reports_missing_pymupdf_clearly(self):
        with self.assertRaisesRegex(RuntimeError, "需要安装 PyMuPDF：pip install pymupdf"):
            build_pdf_project("missing.pdf", mode="quant")


if __name__ == "__main__":
    unittest.main()
