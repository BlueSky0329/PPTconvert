import unittest

from pptx import Presentation
from pptx.util import Inches

from core.template_style import extract_style_from_slide


class TemplateStyleTest(unittest.TestCase):
    def test_extract_style_from_slide_uses_tagged_boxes(self):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])

        slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1)).text = "[stem]"
        slide.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(0.7)).text = "A. option"
        slide.shapes.add_textbox(Inches(5), Inches(2), Inches(3), Inches(0.7)).text = "B. option"
        slide.shapes.add_textbox(Inches(1), Inches(3), Inches(3), Inches(0.7)).text = "C. option"
        slide.shapes.add_textbox(Inches(5), Inches(3), Inches(3), Inches(0.7)).text = "D. option"

        style = extract_style_from_slide(slide)

        self.assertIsNotNone(style.stem_rect)
        self.assertEqual(len(style.option_rects), 4)
        self.assertEqual(len(style.option_box_styles), 4)
