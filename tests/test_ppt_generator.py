import unittest

from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from core.models import Option, Question
from core.ppt_generator import PPTConfig, PPTGenerator


class PPTConfigTest(unittest.TestCase):
    def test_from_mapping_applies_overrides(self):
        config = PPTConfig.from_mapping(
            {
                "margin_left_in": 1.2,
                "stem_font_size": Pt(24),
                "image_max_width": Inches(4),
                "option_letter_color": RGBColor(1, 2, 3),
            }
        )
        self.assertEqual(config.margin_left_in, 1.2)
        self.assertEqual(config.stem_font_size.pt, 24)
        self.assertEqual(config.image_max_width.inches, 4)
        self.assertEqual(config.option_letter_color, RGBColor(1, 2, 3))
        self.assertEqual(config.number_color, RGBColor(1, 2, 3))


class PPTGeneratorTest(unittest.TestCase):
    def test_renderable_options_warns_when_question_has_more_than_four_options(self):
        question = Question(
            number=1,
            stem="stem",
            options=[
                Option("A", "a"),
                Option("B", "b"),
                Option("C", "c"),
                Option("D", "d"),
                Option("E", "e"),
            ],
        )
        generator = PPTGenerator()

        with self.assertLogs("core.ppt_generator", level="WARNING") as captured:
            options = generator._renderable_options(question)

        self.assertEqual([option.letter for option in options], ["A", "B", "C", "D"])
        self.assertIn("only the first 4 can be rendered", captured.output[0])

    def test_stem_text_for_data_question_uses_material_then_source_number(self):
        question = Question(
            number=99,
            stem="2024年全国地表水中未入渗补给地下水的资源量比入渗补给地下水的多：",
            source_question_number="1",
            material_header="材料一",
            material_text="2024年，全国平均年降水量为717.7毫米。",
        )
        text = PPTGenerator._stem_text_for_question(question)
        self.assertEqual(
            text,
            "材料一\n2024年，全国平均年降水量为717.7毫米。\n1. 2024年全国地表水中未入渗补给地下水的资源量比入渗补给地下水的多：",
        )

    def test_layout_default_prefers_question_option_layout_over_global_config(self):
        question = Question(
            number=1,
            stem="stem",
            options=[
                Option("A", "1"),
                Option("B", "2"),
                Option("C", "3"),
                Option("D", "4"),
            ],
            option_layout="one_row",
        )
        generator = PPTGenerator(config=PPTConfig(option_layout="list"))
        generator._prs = generator.tm.create_default()
        slide = generator._prs.slides.add_slide(generator.tm.get_blank_layout())
        called: list[str] = []

        def mark_one_row(*args, **kwargs):
            called.append("one_row")

        def mark_list(*args, **kwargs):
            called.append("list")

        def mark_grid(*args, **kwargs):
            called.append("grid")

        generator._options_one_row = mark_one_row  # type: ignore[method-assign]
        generator._options_list = mark_list  # type: ignore[method-assign]
        generator._options_grid = mark_grid  # type: ignore[method-assign]

        generator._layout_default(slide, question)

        self.assertEqual(called, ["one_row"])

    def test_layout_default_honors_question_ppt_layout_override(self):
        question = Question(
            number=1,
            stem="stem",
            options=[
                Option("A", "1"),
                Option("B", "2"),
                Option("C", "3"),
                Option("D", "4"),
            ],
            ppt_layout={
                "stem": {"x": 0.12, "y": 0.10, "w": 0.70, "h": 0.20},
                "options": {"x": 0.10, "y": 0.52, "w": 0.78, "h": 0.24},
            },
        )
        generator = PPTGenerator()
        generator._prs = generator.tm.create_default()
        slide = generator._prs.slides.add_slide(generator.tm.get_blank_layout())
        calls: dict[str, tuple[int, int, int, int]] = {}

        def capture_stem(_slide, _question, left, top, width, height, style_override=None):
            calls["stem"] = (left, top, width, height)

        def capture_options(_slide, _question, effective_layout):
            rect = effective_layout["options"]
            slide_width = generator._prs.slide_width
            slide_height = generator._prs.slide_height
            calls["options"] = (
                int(slide_width * rect["x"]),
                int(slide_height * rect["y"]),
                int(slide_width * rect["w"]),
                int(slide_height * rect["h"]),
            )

        generator._add_stem_box = capture_stem  # type: ignore[method-assign]
        generator._layout_options_in_region = capture_options  # type: ignore[method-assign]

        generator._layout_default(slide, question)

        slide_width = generator._prs.slide_width
        slide_height = generator._prs.slide_height
        self.assertEqual(
            calls["stem"],
            (
                int(slide_width * 0.12),
                int(slide_height * 0.10),
                int(slide_width * 0.70),
                int(slide_height * 0.20),
            ),
        )
        self.assertEqual(
            calls["options"],
            (
                int(slide_width * 0.10),
                int(slide_height * 0.52),
                int(slide_width * 0.78),
                int(slide_height * 0.24),
            ),
        )

    def test_layout_default_honors_per_option_ppt_layout_override(self):
        question = Question(
            number=1,
            stem="stem",
            options=[
                Option("A", "1"),
                Option("B", "2"),
                Option("C", "3"),
                Option("D", "4"),
            ],
            ppt_layout={
                "options": {"x": 0.10, "y": 0.50, "w": 0.80, "h": 0.28},
                "option_a": {"x": 0.14, "y": 0.54, "w": 0.24, "h": 0.10},
            },
        )
        generator = PPTGenerator()
        generator._prs = generator.tm.create_default()
        slide = generator._prs.slides.add_slide(generator.tm.get_blank_layout())
        calls: dict[str, tuple[int, int, int, int]] = {}

        def capture_option(_slide, letter, _text, left, top, width, height, letter_tpl=None, body_tpl=None):
            calls[letter] = (left, top, width, height)

        generator._add_option_box = capture_option  # type: ignore[method-assign]

        generator._layout_default(slide, question)

        slide_width = generator._prs.slide_width
        slide_height = generator._prs.slide_height
        expected = (
            int(slide_width * 0.14),
            int(slide_height * 0.54),
            int(slide_width * 0.24),
            int(slide_height * 0.10),
        )
        for actual_value, expected_value in zip(calls["A"], expected):
            self.assertLessEqual(abs(actual_value - expected_value), 1)
