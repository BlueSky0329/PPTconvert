"""
从模板 PPT 的首页提取字体、颜色、对齐及文本框位置，用于生成题目页时套用。

支持通过「占位文字」或形状名称标记区域，便于精确定位题干 / 图片区 / 选项：
- 题干：[题干] 或 [stem]；或形状名含 题干 / stem / question；或「标题」占位符
- 图片：插入的图片占位；或文本框写 [图片] / [image]；或形状名含 图片 / pic
- 选项：[选项A]～[选项D]；或 A. / A． 开头的短标签行

组合（Group）内的文本框会递归展开，并使用在幻灯片上的绝对坐标。
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from typing import TYPE_CHECKING, Iterator, Optional, Any

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

if TYPE_CHECKING:
    from pptx.slide import Slide
    from core.ppt_generator import PPTConfig


@dataclass
class TextBoxRect:
    """文本框在幻灯片上的位置（EMU，与 python-pptx 一致）"""
    left: int
    top: int
    width: int
    height: int


@dataclass
class ExtractedRunStyle:
    name: Optional[str] = None
    size_pt: Optional[float] = None
    bold: Optional[bool] = None
    rgb: Optional[RGBColor] = None
    alignment: Optional[int] = None  # PP_ALIGN


@dataclass
class OptionBoxStyle:
    """单个选项框内：字母与正文可不同样式（对应模板中第一、第二 run）"""
    letter: Optional[ExtractedRunStyle] = None
    body: Optional[ExtractedRunStyle] = None


@dataclass
class TemplateSlideStyle:
    """从模板首页解析出的样式（可部分为空，生成时与 GUI 配置合并）"""
    stem: Optional[ExtractedRunStyle] = None
    option: Optional[ExtractedRunStyle] = None  # 兼容：首个选项框整体样式
    stem_rect: Optional[TextBoxRect] = None
    image_rect: Optional[TextBoxRect] = None
    option_rects: list[TextBoxRect] = field(default_factory=list)
    option_box_styles: list[OptionBoxStyle] = field(default_factory=list)
    source_slide_index: int = 0


_RE_OPT_TAG = re.compile(r"\[\s*选项\s*([ABCDabcd])\s*\]")
_RE_OPT_LINE = re.compile(r"^\s*([ABCDabcd])\s*[．.、:：\)）]\s*")
# 一行内多个选项：A. xxx B. xxx …
_RE_OPT_MARKERS = re.compile(r"[ABCDabcd]\s*[\.．、]")

# 主题色无法从 XML 快速解析时，用 Office 默认主题的近似 RGB（仅作兜底）
_THEME_FALLBACK_RGB: dict[MSO_THEME_COLOR, RGBColor] = {
    MSO_THEME_COLOR.TEXT_1: RGBColor(0, 0, 0),
    MSO_THEME_COLOR.TEXT_2: RGBColor(68, 84, 106),
    MSO_THEME_COLOR.DARK_1: RGBColor(0, 0, 0),
    MSO_THEME_COLOR.DARK_2: RGBColor(68, 84, 106),
    MSO_THEME_COLOR.LIGHT_1: RGBColor(255, 255, 255),
    MSO_THEME_COLOR.LIGHT_2: RGBColor(231, 230, 230),
    MSO_THEME_COLOR.BACKGROUND_1: RGBColor(255, 255, 255),
    MSO_THEME_COLOR.BACKGROUND_2: RGBColor(231, 230, 230),
    MSO_THEME_COLOR.ACCENT_1: RGBColor(68, 114, 196),
    MSO_THEME_COLOR.ACCENT_2: RGBColor(237, 125, 49),
    MSO_THEME_COLOR.ACCENT_3: RGBColor(165, 165, 165),
    MSO_THEME_COLOR.ACCENT_4: RGBColor(255, 192, 0),
    MSO_THEME_COLOR.ACCENT_5: RGBColor(91, 155, 213),
    MSO_THEME_COLOR.ACCENT_6: RGBColor(112, 173, 71),
    MSO_THEME_COLOR.HYPERLINK: RGBColor(5, 99, 193),
    MSO_THEME_COLOR.FOLLOWED_HYPERLINK: RGBColor(149, 79, 114),
}


def _rgb_from_font_enhanced(font) -> Optional[RGBColor]:
    """RGB / 主题色 / 部分环境下可直接解析的 rgb。"""
    try:
        if font.color.type == MSO_COLOR_TYPE.RGB and font.color.rgb:
            return font.color.rgb
    except Exception:
        pass
    try:
        rgb = font.color.rgb
        if rgb is not None:
            return rgb
    except Exception:
        pass
    try:
        if font.color.type == MSO_COLOR_TYPE.SCHEME:
            tc = font.color.theme_color
            if tc in _THEME_FALLBACK_RGB:
                return _THEME_FALLBACK_RGB[tc]
    except Exception:
        pass
    return None


def _run_style_from_run(run, paragraph=None) -> Optional[ExtractedRunStyle]:
    if run is None:
        return None
    f = run.font
    size_pt = None
    try:
        if f.size:
            size_pt = float(f.size.pt)
    except Exception:
        pass
    bold = f.bold if f.bold is not None else False
    name = f.name
    rgb = _rgb_from_font_enhanced(f)
    align = PP_ALIGN.LEFT
    par = paragraph
    if par is None:
        try:
            par = getattr(run, "_parent", None) or getattr(run, "parent", None)
        except Exception:
            par = None
    try:
        if par is not None and par.alignment is not None:
            align = par.alignment
    except Exception:
        pass
    return ExtractedRunStyle(
        name=name,
        size_pt=size_pt,
        bold=bold,
        rgb=rgb,
        alignment=align,
    )


def _extract_stem_style(shape) -> Optional[ExtractedRunStyle]:
    """题干：优先取首段中第一个非空 run；否则首 run。"""
    if not shape.has_text_frame:
        return None
    tf = shape.text_frame
    for p in tf.paragraphs:
        for r in p.runs:
            if (r.text or "").strip():
                return _run_style_from_run(r, p)
    for p in tf.paragraphs:
        if p.runs:
            return _run_style_from_run(p.runs[0], p)
    return None


def _option_box_style_from_shape(shape) -> OptionBoxStyle:
    """选项：若有两段以上 run，分别作为「字母」「正文」样式；否则共用。"""
    if not shape.has_text_frame:
        return OptionBoxStyle()
    p = shape.text_frame.paragraphs[0]
    if not p.runs:
        st = _extract_stem_style(shape)
        return OptionBoxStyle(st, st)
    if len(p.runs) >= 2:
        a = _run_style_from_run(p.runs[0], p)
        b = _run_style_from_run(p.runs[1], p)
        return OptionBoxStyle(a, b or a)
    st = _run_style_from_run(p.runs[0], p)
    return OptionBoxStyle(st, st)


def _rect_from_shape(shape, abs_left: int, abs_top: int) -> TextBoxRect:
    return TextBoxRect(
        left=int(abs_left),
        top=int(abs_top),
        width=int(shape.width),
        height=int(shape.height),
    )


def _iter_slide_shapes_flat(slide: "Slide") -> Iterator[tuple[object, int, int]]:
    """
    遍历幻灯片上所有形状（含组合内），返回 (shape, 绝对 left, 绝对 top)。
    子形状在组合内为相对坐标，此处累加为幻灯片坐标。
    """
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_group(sh, int(sh.left), int(sh.top))
        else:
            yield sh, int(sh.left), int(sh.top)


def _walk_group(group, bx: int, by: int) -> Iterator[tuple[object, int, int]]:
    for sh in group.shapes:
        al = bx + int(sh.left)
        at = by + int(sh.top)
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_group(sh, al, at)
        else:
            yield sh, al, at


def _sorted_flat_text_shapes(slide: "Slide"):
    """含组合内文本框，按绝对位置排序。"""
    out = []
    for sh, al, at in _iter_slide_shapes_flat(slide):
        if getattr(sh, "has_text_frame", False) and sh.has_text_frame:
            out.append((sh, al, at))
    out.sort(key=lambda t: (t[2], t[1]))
    return out


def _shape_plain_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
        return ""
    return (shape.text_frame.text or "").strip()


def _classify_shape(shape) -> tuple[str, Optional[int]]:
    """
    返回 (角色, 选项下标 0~3)。
    角色：stem | image | image_rect | option | unknown
    """
    try:
        st = shape.shape_type
        if st == MSO_SHAPE_TYPE.PICTURE:
            return ("image", None)
    except Exception:
        pass

    try:
        if getattr(shape, "is_placeholder", False) and shape.is_placeholder:
            ph = shape.placeholder_format
            if ph.type in (
                PP_PLACEHOLDER_TYPE.PICTURE,
                PP_PLACEHOLDER_TYPE.BITMAP,
            ):
                return ("image", None)
            if ph.type in (
                PP_PLACEHOLDER_TYPE.TITLE,
                PP_PLACEHOLDER_TYPE.CENTER_TITLE,
                PP_PLACEHOLDER_TYPE.VERTICAL_TITLE,
            ):
                return ("stem", None)
    except Exception:
        pass

    name = (getattr(shape, "name", "") or "").strip().lower()
    txt = _shape_plain_text(shape)
    txt_low = txt.lower()

    if "[题干]" in txt or "[stem]" in txt_low or "【题干】" in txt:
        return ("stem", None)
    if "[图片]" in txt or "[image]" in txt_low or "【图片】" in txt:
        return ("image_rect", None)

    m = _RE_OPT_TAG.search(txt)
    if m:
        idx = ord(m.group(1).upper()) - ord("A")
        if 0 <= idx <= 3:
            return ("option", idx)

    if txt and len(txt) < 120:
        m2 = _RE_OPT_LINE.match(txt)
        if m2:
            idx = ord(m2.group(1).upper()) - ord("A")
            if 0 <= idx <= 3:
                return ("option", idx)

    if any(k in name for k in ("题干", "stem", "question", "ti_mu", "timu")):
        return ("stem", None)
    if any(k in name for k in ("图片", "image", "pic", "figure", "chart", "photo")):
        return ("image_rect", None)
    for i, ch in enumerate("abcd"):
        if f"option_{ch}" in name or f"opt_{ch}" in name or f"选项{chr(65+i).lower()}" in name:
            return ("option", i)
        if f"选项{chr(65+i)}" in name:
            return ("option", i)

    return ("unknown", None)


def extract_style_from_slide(slide: "Slide") -> TemplateSlideStyle:
    """
    从单页幻灯片推断样式与矩形：
    1) 优先根据 [题干]/[图片]/[选项A]～、标题占位符、形状名识别；
    2) 否则按文本框绝对位置：第 1 个为题干，其后顺次填满 A～D（须凑齐 4 个选项框）。
    """
    style = TemplateSlideStyle()
    all_shapes = []
    for sh, al, at in _iter_slide_shapes_flat(slide):
        all_shapes.append((sh, al, at))

    classified = [(sh, al, at, _classify_shape(sh)) for sh, al, at in all_shapes]

    stem_sh = None
    stem_abs = (0, 0)
    image_rect: Optional[TextBoxRect] = None
    opt_shapes: list[Optional[object]] = [None, None, None, None]
    opt_abs: list[tuple[int, int]] = [(0, 0), (0, 0), (0, 0), (0, 0)]

    for sh, al, at, (role, idx) in classified:
        if role == "stem" and stem_sh is None:
            stem_sh = sh
            stem_abs = (al, at)
        elif role == "image":
            image_rect = _rect_from_shape(sh, al, at)
        elif role == "image_rect":
            image_rect = _rect_from_shape(sh, al, at)
        elif role == "option" and idx is not None:
            opt_shapes[idx] = sh
            opt_abs[idx] = (al, at)

    text_shapes = _sorted_flat_text_shapes(slide)
    layout_text_shapes = []
    for sh, al, at in text_shapes:
        r, _ = _classify_shape(sh)
        if r == "image_rect":
            if image_rect is None:
                image_rect = _rect_from_shape(sh, al, at)
            continue
        layout_text_shapes.append((sh, al, at))

    used_ids = set()
    if stem_sh is not None:
        used_ids.add(id(stem_sh))
    for sh in opt_shapes:
        if sh is not None:
            used_ids.add(id(sh))

    remaining = [(s, al, at) for s, al, at in layout_text_shapes if id(s) not in used_ids]

    if stem_sh is None and remaining:
        stem_sh, sx, sy = remaining[0]
        stem_abs = (sx, sy)
        used_ids.add(id(stem_sh))
        remaining = [(s, al, at) for s, al, at in remaining[1:] if id(s) not in used_ids]

    for i in range(4):
        if opt_shapes[i] is None and remaining:
            s, al, at = remaining.pop(0)
            opt_shapes[i] = s
            opt_abs[i] = (al, at)

    if stem_sh is not None:
        style.stem = _extract_stem_style(stem_sh)
        style.stem_rect = _rect_from_shape(stem_sh, stem_abs[0], stem_abs[1])

    if all(opt_shapes):
        style.option_rects = [_rect_from_shape(opt_shapes[i], opt_abs[i][0], opt_abs[i][1]) for i in range(4)]
        style.option_box_styles = [_option_box_style_from_shape(opt_shapes[i]) for i in range(4)]
        style.option = _extract_stem_style(opt_shapes[0])
    else:
        for sh in opt_shapes:
            if sh is not None and style.option is None:
                style.option = _extract_stem_style(sh)
                break

    if image_rect is None:
        for sh, al, at, (role, _) in classified:
            if role == "image":
                image_rect = _rect_from_shape(sh, al, at)
                break

    style.image_rect = image_rect

    _apply_combined_option_box(style, slide, stem_sh, opt_shapes, opt_abs)
    return style


def neutralize_option_colors_if_no_template_rgb(
    cfg: "PPTConfig", template_style: TemplateSlideStyle
) -> None:
    """
    模板合并后：若聚合的选项样式里没有解析到 RGB（主题色未落盘等），
    则不要用 PPTConfig 类默认的「选项字母蓝」#006BBD，改为黑色，避免与模板实际黑字不符。
    """
    o = template_style.option
    if o is None or o.rgb is None:
        z = RGBColor(0, 0, 0)
        cfg.option_letter_color = z
        cfg.option_color = z
        cfg.number_color = z


def merge_template_style_into_config(cfg: "PPTConfig", template_style: TemplateSlideStyle) -> None:
    """用模板样式覆盖 config（题干、选项分别覆盖；字体名优先采用选项框字体）"""
    if template_style.stem:
        s = template_style.stem
        if s.name:
            cfg.font_name = s.name
        if s.size_pt is not None and s.size_pt > 0:
            cfg.stem_font_size = Pt(s.size_pt)
        if s.bold is not None:
            cfg.font_bold_stem = bool(s.bold)
        if s.rgb is not None:
            cfg.stem_color = s.rgb
        if s.alignment is not None:
            cfg.stem_align = _align_to_str(s.alignment)

    if template_style.option:
        o = template_style.option
        if o.name:
            cfg.font_name = o.name
        if o.size_pt is not None and o.size_pt > 0:
            cfg.option_font_size = Pt(o.size_pt)
        if o.bold is not None:
            cfg.option_letter_bold = bool(o.bold)
            cfg.option_font_bold = bool(o.bold)
        if o.rgb is not None:
            cfg.option_letter_color = o.rgb
            cfg.option_color = o.rgb
            cfg.number_color = o.rgb
        if o.alignment is not None:
            cfg.option_align = _align_to_str(o.alignment)


def _align_to_str(align: int) -> str:
    if align == PP_ALIGN.CENTER:
        return "center"
    if align == PP_ALIGN.RIGHT:
        return "right"
    return "left"


def _looks_like_combined_four_options(txt: str) -> bool:
    """判断是否为「单行/单框」内排布的 A. B. C. D. 四选项。"""
    if not txt or len(txt) < 8:
        return False
    markers = _RE_OPT_MARKERS.findall(txt)
    letters = set()
    for m in markers:
        letters.add(m.strip()[0].upper())
    return letters >= {"A", "B", "C", "D"}


def _quarter_rects_ab_cd(r: TextBoxRect) -> list[TextBoxRect]:
    """将矩形四等分为 2×2：上排 A B，下排 C D（与默认网格一致）。"""
    w2 = max(1, r.width // 2)
    h2 = max(1, r.height // 2)
    return [
        TextBoxRect(r.left, r.top, w2, h2),
        TextBoxRect(r.left + w2, r.top, w2, h2),
        TextBoxRect(r.left, r.top + h2, w2, h2),
        TextBoxRect(r.left + w2, r.top + h2, w2, h2),
    ]


def _apply_combined_option_box(
    style: TemplateSlideStyle,
    slide: "Slide",
    stem_sh: Optional[Any],
    opt_shapes: list[Optional[Any]],
    opt_abs: list[tuple[int, int]],
) -> None:
    """
    若四个选项挤在同一文本框（常见于一行四选项），将原文本框矩形四等分，
    仍沿用该框解析出的字体样式。
    """
    if all(opt_shapes):
        return

    cand_sh = None
    cand_al = cand_at = 0
    for i, sh in enumerate(opt_shapes):
        if sh is not None:
            txt = _shape_plain_text(sh)
            if _looks_like_combined_four_options(txt):
                cand_sh = sh
                cand_al, cand_at = opt_abs[i]
                break

    if cand_sh is None:
        for sh, al, at in _sorted_flat_text_shapes(slide):
            if stem_sh is not None and id(sh) == id(stem_sh):
                continue
            txt = _shape_plain_text(sh)
            if _looks_like_combined_four_options(txt):
                cand_sh = sh
                cand_al, cand_at = al, at
                break

    if cand_sh is None:
        return

    r = _rect_from_shape(cand_sh, cand_al, cand_at)
    style.option_rects = _quarter_rects_ab_cd(r)
    ob = _option_box_style_from_shape(cand_sh)
    style.option_box_styles = [ob, ob, ob, ob]
    if style.option is None:
        style.option = _extract_stem_style(cand_sh)


def delete_all_slides(prs) -> None:
    """删除演示文稿中全部幻灯片，保留母版与版式，仅清空 sldIdLst。"""
    sld_id_lst = prs.slides._sldIdLst
    ids = list(sld_id_lst.sldId_lst)
    for sld_id in reversed(ids):
        r_id = sld_id.rId
        prs.part.drop_rel(r_id)
        sld_id.getparent().remove(sld_id)
