import logging
import os
from dataclasses import dataclass, field, replace
from typing import Mapping, Optional

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from core.models import Question
from core.ppt_style import align_from_string
from core.template_manager import TemplateManager
from core.template_style import (
    ExtractedRunStyle,
    TemplateSlideStyle,
    delete_all_slides,
    extract_style_from_slide,
    merge_template_style_into_config,
    neutralize_option_colors_if_no_template_rgb,
)

EMU_PER_INCH = 914400
LOGGER = logging.getLogger(__name__)


@dataclass
class PPTConfig:
    """PPT generation settings."""

    margin_left_in: float = 0.8
    margin_right_in: float = 0.8
    margin_top_in: float = 0.5

    stem_height_with_image_in: float = 1.5
    stem_height_no_image_in: float = 2.5
    stem_font_size: object = field(default_factory=lambda: Pt(20))
    font_name: str = "微软雅黑"
    font_bold_stem: bool = True
    stem_color: RGBColor = field(default_factory=lambda: RGBColor(0x1A, 0x1A, 0x2E))
    stem_align: str = "left"

    image_max_width: object = field(default_factory=lambda: Inches(5))
    image_max_height: object = field(default_factory=lambda: Inches(2.5))
    image_h_align: str = "center"

    gap_after_stem_in: float = 0.2
    gap_after_image_in: float = 0.15
    gap_before_options_in: float = 0.2

    option_layout: str = "grid"
    one_row_height_in: float = 0.55
    one_row_gap_in: float = 0.06
    grid_layout: str = "ab_cd"
    grid_row_height_in: float = 0.9
    grid_col_gap_in: float = 0.15
    list_row_height_in: float = 0.7
    option_font_size: object = field(default_factory=lambda: Pt(18))
    option_font_bold: bool = False
    option_color: RGBColor = field(default_factory=lambda: RGBColor(0x2D, 0x2D, 0x2D))
    option_letter_color: RGBColor = field(default_factory=lambda: RGBColor(0x00, 0x6B, 0xBD))
    option_letter_bold: bool = True
    option_align: str = "left"

    number_color: Optional[RGBColor] = None

    def __post_init__(self):
        if self.number_color is None:
            self.number_color = self.option_letter_color

    @classmethod
    def from_mapping(cls, values: Mapping[str, object]) -> "PPTConfig":
        config = cls()
        for key, value in values.items():
            if hasattr(config, key) and value is not None:
                setattr(config, key, value)
        config.number_color = config.option_letter_color
        return config

    def stem_align_pp(self) -> int:
        return align_from_string(self.stem_align)

    def option_align_pp(self) -> int:
        return align_from_string(self.option_align)


def _scale_image(img_path: str, max_w_in: float, max_h_in: float):
    px_w = px_h = None
    dpi = 96.0
    try:
        with Image.open(img_path) as img:
            px_w, px_h = img.size
    except Exception:
        LOGGER.warning(
            "Failed to inspect image size for %s; using fallback size",
            img_path,
            exc_info=True,
        )

    if px_w and px_h and px_w > 0 and px_h > 0:
        img_w = px_w / dpi
        img_h = px_h / dpi
        scale = min(max_w_in / img_w, max_h_in / img_h, 1.0)
        return Inches(img_w * scale), Inches(img_h * scale)
    return Inches(min(max_w_in, 4.0)), Inches(min(max_h_in, 2.5))


def _align_pic_left(base_left: int, area_width: int, pic_width, h_align: str) -> int:
    picture_width = int(pic_width)
    if h_align == "right":
        return base_left + area_width - picture_width
    if h_align == "center":
        return max(base_left, base_left + (area_width - picture_width) // 2)
    return base_left


class PPTGenerator:
    def __init__(
        self,
        template_manager: Optional[TemplateManager] = None,
        config: Optional[PPTConfig] = None,
    ):
        self.tm = template_manager or TemplateManager()
        self.config = config or PPTConfig()
        self._prs: Optional[Presentation] = None
        self._tpl_style: Optional[TemplateSlideStyle] = None

    def generate(
        self,
        questions: list[Question],
        output_path: str,
        template_path: Optional[str] = None,
        progress_callback=None,
    ):
        self._tpl_style = None

        if template_path:
            self._prs = self.tm.load_template(template_path)
            if len(self._prs.slides) > 0:
                first = next(iter(self._prs.slides))
                self._tpl_style = extract_style_from_slide(first)
            self.config = replace(self.config)
            if self._tpl_style:
                merge_template_style_into_config(self.config, self._tpl_style)
                neutralize_option_colors_if_no_template_rgb(self.config, self._tpl_style)
            delete_all_slides(self._prs)
        else:
            self._prs = self.tm.create_default()

        total = len(questions)
        for idx, question in enumerate(questions):
            self._add_question_slide(question)
            if progress_callback:
                progress_callback(idx + 1, total)

        output_dir = os.path.dirname(os.path.abspath(output_path))
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)
        self._prs.save(output_path)
        return output_path

    def _renderable_options(self, question: Question) -> list:
        if len(question.options) > 4:
            LOGGER.warning(
                "Question %s has %s options; only the first 4 can be rendered",
                question.number,
                len(question.options),
            )
        return question.options[:4]

    @staticmethod
    def _stem_text_for_question(question: Question) -> str:
        display_number = (question.source_question_number or "").strip() or str(question.number)
        question_line = f"{display_number}. {question.stem}".strip() if question.stem else f"{display_number}."
        if question.material_header or question.material_text:
            parts = [
                part
                for part in (question.material_header, question.material_text, question_line)
                if part
            ]
            return "\n".join(parts)
        display_stem = question.display_stem
        return f"{display_number}. {display_stem}".strip() if display_stem else f"{display_number}."

    def _add_question_slide(self, question: Question):
        slide = self._prs.slides.add_slide(self.tm.get_blank_layout())
        template_style = self._tpl_style
        if template_style and template_style.stem_rect and len(template_style.option_rects) >= 4:
            self._layout_from_template(slide, question, template_style)
        else:
            self._layout_default(slide, question)

    def _layout_from_template(self, slide, question: Question, template_style: TemplateSlideStyle):
        config = self.config
        stem_rect = template_style.stem_rect
        render_options = self._renderable_options(question)

        self._add_stem_box(
            slide,
            question,
            stem_rect.left,
            stem_rect.top,
            stem_rect.width,
            stem_rect.height,
            style_override=template_style.stem,
        )

        stem_bottom = stem_rect.top + stem_rect.height
        self._place_images_template(slide, question, template_style, config, stem_bottom)

        for idx, option in enumerate(render_options):
            if idx >= len(template_style.option_rects):
                break
            rect = template_style.option_rects[idx]
            letter_style, body_style = self._resolve_option_style(template_style, idx)
            self._add_option_box(
                slide,
                option.letter,
                option.text,
                rect.left,
                rect.top,
                rect.width,
                rect.height,
                letter_tpl=letter_style,
                body_tpl=body_style,
            )

    def _place_images_template(self, slide, question, template_style, config, stem_bottom):
        if not question.image_paths:
            return
        if template_style.image_rect:
            image_rect = template_style.image_rect
            top = image_rect.top
            for path in question.image_paths:
                if os.path.exists(path):
                    top = self._insert_image_in_rect(
                        slide,
                        path,
                        image_rect.left,
                        image_rect.top,
                        image_rect.width,
                        image_rect.height,
                        top,
                    )
        else:
            gap_after = int(Inches(config.gap_after_stem_in))
            gap_before = int(Inches(config.gap_before_options_in))
            option_top = min(rect.top for rect in template_style.option_rects)
            top = stem_bottom + gap_after
            band_height = max(1, option_top - top - gap_before)
            for path in question.image_paths:
                if os.path.exists(path):
                    top = self._insert_image(
                        slide,
                        path,
                        template_style.stem_rect.left,
                        top,
                        template_style.stem_rect.width,
                        band_height,
                    )

    @staticmethod
    def _resolve_option_style(template_style: TemplateSlideStyle, idx: int):
        letter_style = body_style = None
        if template_style.option_box_styles and idx < len(template_style.option_box_styles):
            option_box = template_style.option_box_styles[idx]
            letter_style, body_style = option_box.letter, option_box.body
        if letter_style is None and body_style is None and template_style.option:
            letter_style = body_style = template_style.option
        return letter_style, body_style

    def _layout_default(self, slide, question: Question):
        config = self.config
        slide_width = self._prs.slide_width
        margin_left = Inches(config.margin_left_in)
        margin_right = Inches(config.margin_right_in)
        margin_top = Inches(config.margin_top_in)
        content_width = slide_width - margin_left - margin_right
        has_image = bool(question.image_paths)

        stem_height = Inches(
            config.stem_height_with_image_in if has_image else config.stem_height_no_image_in
        )
        self._add_stem_box(slide, question, margin_left, margin_top, content_width, stem_height)

        top = margin_top + stem_height + Inches(config.gap_after_stem_in)
        if has_image:
            for path in question.image_paths:
                if os.path.exists(path):
                    top = self._insert_image(slide, path, margin_left, top, content_width)

        options_top = top + Inches(config.gap_before_options_in)
        option_layout = (config.option_layout or "grid").lower()
        if option_layout == "list":
            self._options_list(slide, question, margin_left, options_top, content_width)
        elif option_layout == "one_row":
            self._options_one_row(slide, question, margin_left, options_top, content_width)
        else:
            self._options_grid(slide, question, margin_left, options_top, content_width)

    def _add_stem_box(
        self,
        slide,
        question: Question,
        left,
        top,
        width,
        height,
        style_override: Optional[ExtractedRunStyle] = None,
    ):
        config = self.config
        box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = box.text_frame
        text_frame.word_wrap = True
        paragraph = text_frame.paragraphs[0]
        paragraph.text = self._stem_text_for_question(question)

        extracted = style_override
        paragraph.font.size = (
            Pt(extracted.size_pt)
            if extracted and extracted.size_pt and extracted.size_pt > 0
            else config.stem_font_size
        )
        paragraph.font.bold = (
            extracted.bold if extracted and extracted.bold is not None else config.font_bold_stem
        )
        paragraph.font.name = extracted.name if extracted and extracted.name else config.font_name
        paragraph.font.color.rgb = extracted.rgb if extracted and extracted.rgb else config.stem_color
        paragraph.alignment = (
            extracted.alignment
            if extracted and extracted.alignment is not None
            else config.stem_align_pp()
        )

    def _insert_image(self, slide, img_path: str, left: int, top: int, area_w: int, max_h_emu: int = 0) -> int:
        config = self.config
        max_width = config.image_max_width.inches
        max_height = config.image_max_height.inches
        if max_h_emu > 0:
            max_height = min(max_height, max_h_emu / EMU_PER_INCH)
        max_width = min(max_width, area_w / EMU_PER_INCH)

        width, height = _scale_image(img_path, max_width, max_height)
        picture_left = _align_pic_left(left, area_w, width, (config.image_h_align or "center").lower())
        picture = slide.shapes.add_picture(img_path, picture_left, top, width=width, height=height)
        return int(top + picture.height + Inches(config.gap_after_image_in))

    def _insert_image_in_rect(
        self,
        slide,
        img_path: str,
        rect_l: int,
        rect_t: int,
        rect_w: int,
        rect_h: int,
        start_top: int,
    ) -> int:
        config = self.config
        bottom = rect_t + rect_h
        remaining = bottom - start_top
        if remaining <= 0:
            return start_top

        max_width = min(config.image_max_width.inches, rect_w / EMU_PER_INCH)
        max_height = min(config.image_max_height.inches, remaining / EMU_PER_INCH)
        width, height = _scale_image(img_path, max_width, max_height)
        picture_left = _align_pic_left(
            rect_l,
            rect_w,
            width,
            (config.image_h_align or "center").lower(),
        )
        picture = slide.shapes.add_picture(img_path, picture_left, start_top, width=width, height=height)
        next_top = int(start_top + picture.height + Inches(config.gap_after_image_in))
        return min(next_top, bottom)

    def _add_option_box(
        self,
        slide,
        letter: str,
        text: str,
        left,
        top,
        width,
        height,
        letter_tpl: Optional[ExtractedRunStyle] = None,
        body_tpl: Optional[ExtractedRunStyle] = None,
    ):
        config = self.config
        box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = box.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = None

        paragraph = text_frame.paragraphs[0]
        letter_run = paragraph.add_run()
        letter_run.text = f"{letter}. "
        self._style_run(
            letter_run,
            config.option_font_size,
            config.option_letter_bold,
            config.option_letter_color,
            letter_tpl,
        )

        body_run = paragraph.add_run()
        body_run.text = text
        self._style_run(
            body_run,
            config.option_font_size,
            config.option_font_bold,
            config.option_color,
            body_tpl or letter_tpl,
        )

        if body_tpl and body_tpl.alignment is not None:
            paragraph.alignment = body_tpl.alignment
        elif letter_tpl and letter_tpl.alignment is not None:
            paragraph.alignment = letter_tpl.alignment
        else:
            paragraph.alignment = config.option_align_pp()

    def _style_run(self, run, default_size, default_bold, default_rgb, tpl: Optional[ExtractedRunStyle]):
        config = self.config
        run.font.size = Pt(tpl.size_pt) if tpl and tpl.size_pt and tpl.size_pt > 0 else default_size
        run.font.bold = tpl.bold if tpl and tpl.bold is not None else default_bold
        run.font.name = tpl.name if tpl and tpl.name else config.font_name
        run.font.color.rgb = tpl.rgb if tpl and tpl.rgb else default_rgb

    def _grid_positions(self, left, top, col_w, row_h):
        config = self.config
        gap = Inches(config.grid_col_gap_in)
        col0, col1 = left, left + col_w + gap
        row0, row1 = top, top + row_h
        if config.grid_layout == "ac_bd":
            return [(col0, row0), (col0, row1), (col1, row0), (col1, row1)]
        return [(col0, row0), (col1, row0), (col0, row1), (col1, row1)]

    def _options_grid(self, slide, question: Question, left, top, content_width):
        config = self.config
        gap = Inches(config.grid_col_gap_in)
        col_width = (content_width - gap) // 2
        row_height = Inches(config.grid_row_height_in)
        positions = self._grid_positions(left, top, col_width, row_height)
        width, height = col_width - Inches(0.3), row_height - Inches(0.1)
        for idx, option in enumerate(self._renderable_options(question)):
            if idx < len(positions):
                self._add_option_box(
                    slide,
                    option.letter,
                    option.text,
                    positions[idx][0],
                    positions[idx][1],
                    width,
                    height,
                )

    def _options_list(self, slide, question: Question, left, top, content_width):
        row_height = Inches(self.config.list_row_height_in)
        for idx, option in enumerate(self._renderable_options(question)):
            self._add_option_box(
                slide,
                option.letter,
                option.text,
                left,
                top + idx * row_height,
                content_width,
                row_height - Inches(0.1),
            )

    def _options_one_row(self, slide, question: Question, left, top, content_width):
        config = self.config
        options = self._renderable_options(question)
        if not options:
            return
        option_count = len(options)
        gap = Inches(config.one_row_gap_in)
        cell_width = (content_width - gap * (option_count - 1)) // option_count if option_count > 1 else content_width
        row_height = Inches(config.one_row_height_in)
        pad = Inches(0.04)
        for idx, option in enumerate(options):
            x = left + idx * (cell_width + gap)
            self._add_option_box(
                slide,
                option.letter,
                option.text,
                x,
                top,
                cell_width - pad,
                row_height - pad,
            )
